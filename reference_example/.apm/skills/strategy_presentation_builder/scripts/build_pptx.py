"""
Build a Strategy Presentation deck that matches KaxaNuk branding.

Programmatic helper used by the `strategy_presentation_builder` skill. The skill
collects strategy content from the project context (memory, backtest outputs,
scoring spreadsheet) and passes it to `build_deck` as a structured `Deck` dict.
The function then renders the deck slide-by-slide using `python-pptx`, applying
the brand templates documented in `references/BRANDING_GUIDELINES.md` and the
slide order documented in `references/SLIDE_STRUCTURE.md`.

Output format is controlled by `--format`:

- `pptx` (default): writes an editable `.pptx`.
- `pdf`: writes the same deck and then converts it to PDF using LibreOffice
  (`soffice --headless --convert-to pdf`). LibreOffice must be on PATH; if it
  is not, the script raises so the caller can install it or pick `pptx`.

Examples:

    python build_pptx.py --example --output /tmp/demo.pptx
    python build_pptx.py --input deck.json --format pdf --output out/deck.pdf

Dependencies: python-pptx (>=0.6.21). For PDF output, LibreOffice ≥ 7.x.
"""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Brand constants — keep in sync with references/BRANDING_GUIDELINES.md.
# ---------------------------------------------------------------------------

BRAND_RED_PRIMARY = RGBColor(0xCB, 0x2E, 0x15)
BRAND_RED_ACCENT = RGBColor(0xE8, 0x43, 0x28)
CHARCOAL = RGBColor(0x2C, 0x2C, 0x2C)
SUBTITLE_GRAY = RGBColor(0xDC, 0xDC, 0xDC)
PAGE_NUM_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINK_BLUE = RGBColor(0x11, 0x55, 0xCC)

FONT_PRIMARY = "Montserrat"
FONT_FALLBACK = "Arial"

# Slide dimensions (16:9 widescreen).
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Asset paths resolved relative to this file.
ASSETS_DIR = Path(__file__).resolve().parent.parent / "references" / "assets"
LOGO_RED = ASSETS_DIR / "kaxanuk_logo_red.png"
WORDMARK = ASSETS_DIR / "kaxanuk_wordmark.png"
FINGERPRINT_GRAY = ASSETS_DIR / "kaxanuk_fingerprint_gray.png"


# ---------------------------------------------------------------------------
# Data model.
# ---------------------------------------------------------------------------

@dataclass
class Slide:
    """One slide entry in the deck.

    Layouts:
        - "cover"    -> wordmark cover (no title)
        - "title"    -> brand-red title slide
        - "divider"  -> charcoal section divider
        - "content"  -> white slide with red title and bullet body
        - "visual"   -> white slide with red title and a single embedded image
        - "contact"  -> final contact slide
    """
    layout: str
    title: str | None = None
    subtitle: str | None = None
    body: list[dict] = field(default_factory=list)
    image_path: str | None = None
    page_number: int | None = None


@dataclass
class Deck:
    strategy_name: str
    contact_email: str = "research@kaxanuk.mx"
    slides: list[Slide] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Low-level helpers.
# ---------------------------------------------------------------------------

def _set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank.


def _fill_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = FONT_PRIMARY,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bulleted_body(
    slide,
    blocks: Iterable[dict],
    *,
    left: float = 0.6,
    top: float = 1.45,
    width: float = 12.1,
    height: float = 5.4,
):
    """Render a body composed of text blocks.

    Each block is a dict with keys:
        kind:     "para" | "heading" | "bullets" | "arrows" | "callout"
        text:     str (for para/heading/callout)
        items:    list[str] (for bullets/arrows)
        bold:     bool (heading/callout default True)
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True

    def new_para():
        nonlocal first
        if first:
            first = False
            return tf.paragraphs[0]
        return tf.add_paragraph()

    for block in blocks:
        kind = block.get("kind", "para")
        if kind in ("para", "heading", "callout"):
            text = block.get("text", "")
            p = new_para()
            run = p.add_run()
            run.text = text
            run.font.name = FONT_PRIMARY
            run.font.size = Pt(16 if kind == "heading" else 14)
            run.font.bold = block.get("bold", kind in ("heading", "callout"))
            run.font.color.rgb = CHARCOAL
            p.space_after = Pt(6)
        elif kind == "bullets":
            for item in block.get("items", []):
                p = new_para()
                run = p.add_run()
                run.text = f"●  {item}"  # filled disc + two spaces
                run.font.name = FONT_PRIMARY
                run.font.size = Pt(14)
                run.font.color.rgb = CHARCOAL
                p.space_after = Pt(2)
        elif kind == "arrows":
            for item in block.get("items", []):
                p = new_para()
                run_arrow = p.add_run()
                run_arrow.text = "➔  "  # heavy right arrow
                run_arrow.font.name = FONT_PRIMARY
                run_arrow.font.size = Pt(14)
                run_arrow.font.bold = True
                run_arrow.font.color.rgb = BRAND_RED_ACCENT
                run_text = p.add_run()
                run_text.text = item
                run_text.font.name = FONT_PRIMARY
                run_text.font.size = Pt(14)
                run_text.font.color.rgb = CHARCOAL
                p.space_after = Pt(4)
        else:
            raise ValueError(f"Unknown body block kind: {kind!r}")
    return box


def _add_logo(slide) -> None:
    if LOGO_RED.exists():
        slide.shapes.add_picture(
            str(LOGO_RED),
            Inches(12.45),
            Inches(0.30),
            width=Inches(0.65),
            height=Inches(0.65),
        )


def _add_page_number(slide, n: int) -> None:
    _add_text_box(
        slide,
        str(n),
        left=12.55,
        top=7.10,
        width=0.6,
        height=0.3,
        size=9,
        color=PAGE_NUM_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def _add_title(slide, text: str) -> None:
    _add_text_box(
        slide,
        text,
        left=0.5,
        top=0.45,
        width=12.0,
        height=0.9,
        size=32,
        color=BRAND_RED_ACCENT,
    )


# ---------------------------------------------------------------------------
# Layout renderers.
# ---------------------------------------------------------------------------

def render_cover(prs: Presentation, slide_def: Slide) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    if FINGERPRINT_GRAY.exists():
        slide.shapes.add_picture(
            str(FINGERPRINT_GRAY),
            Inches(0.4),
            Inches(0.85),
            width=Inches(5.9),
            height=Inches(5.9),
        )
    if WORDMARK.exists():
        slide.shapes.add_picture(
            str(WORDMARK),
            Inches(6.4),
            Inches(2.95),
            width=Inches(6.5),
            height=Inches(1.5),
        )


def render_title(prs: Presentation, slide_def: Slide) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, BRAND_RED_PRIMARY)
    _add_text_box(
        slide,
        slide_def.title or "Strategy Presentation",
        left=0.5,
        top=2.8,
        width=12.3,
        height=1.2,
        size=54,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    if slide_def.subtitle:
        _add_text_box(
            slide,
            slide_def.subtitle,
            left=0.5,
            top=4.05,
            width=12.3,
            height=0.6,
            size=24,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )


def render_divider(prs: Presentation, slide_def: Slide) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, CHARCOAL)
    _add_text_box(
        slide,
        slide_def.title or "",
        left=0.7,
        top=2.95,
        width=11.5,
        height=1.2,
        size=54,
        color=WHITE,
    )
    _add_text_box(
        slide,
        slide_def.subtitle or "",
        left=0.7,
        top=4.20,
        width=11.5,
        height=0.6,
        size=22,
        color=SUBTITLE_GRAY,
    )


def render_content(prs: Presentation, slide_def: Slide) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _add_title(slide, slide_def.title or "")
    _add_logo(slide)
    if slide_def.body:
        _add_bulleted_body(slide, slide_def.body)
    if slide_def.page_number is not None:
        _add_page_number(slide, slide_def.page_number)


def render_visual(prs: Presentation, slide_def: Slide) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    _add_title(slide, slide_def.title or "")
    _add_logo(slide)
    img = slide_def.image_path
    if img and Path(img).exists():
        slide.shapes.add_picture(
            img,
            Inches(1.4),
            Inches(1.6),
            width=Inches(10.5),
            height=Inches(5.0),
        )
    else:
        # Placeholder rectangle so the structure is preserved even when the
        # asset is missing.
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.4),
            Inches(1.6),
            Inches(10.5),
            Inches(5.0),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        rect.line.color.rgb = PAGE_NUM_GRAY
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        label = slide_def.title or "visual"
        run.text = f"<{label} unavailable>"
        run.font.name = FONT_PRIMARY
        run.font.size = Pt(16)
        run.font.color.rgb = PAGE_NUM_GRAY
    if slide_def.page_number is not None:
        _add_page_number(slide, slide_def.page_number)


def render_contact(prs: Presentation, slide_def: Slide, contact_email: str) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, WHITE)
    if FINGERPRINT_GRAY.exists():
        slide.shapes.add_picture(
            str(FINGERPRINT_GRAY),
            Inches(7.0),
            Inches(0.85),
            width=Inches(5.9),
            height=Inches(5.9),
        )
    if LOGO_RED.exists():
        slide.shapes.add_picture(
            str(LOGO_RED),
            Inches(9.3),
            Inches(2.95),
            width=Inches(1.5),
            height=Inches(1.5),
        )
    _add_text_box(
        slide,
        slide_def.title or "Contact",
        left=1.6,
        top=3.10,
        width=4.5,
        height=0.7,
        size=32,
        color=BRAND_RED_ACCENT,
    )
    _add_text_box(
        slide,
        contact_email,
        left=1.6,
        top=3.85,
        width=4.5,
        height=0.5,
        size=18,
        color=LINK_BLUE,
    )


# ---------------------------------------------------------------------------
# Public API.
# ---------------------------------------------------------------------------

LAYOUT_RENDERERS = {
    "cover": render_cover,
    "title": render_title,
    "divider": render_divider,
    "content": render_content,
    "visual": render_visual,
}


def _render_pptx(deck: Deck, output_path: Path) -> Path:
    """Render a Deck to a .pptx file at output_path."""
    prs = Presentation()
    _set_slide_size(prs)
    for s in deck.slides:
        if s.layout == "contact":
            render_contact(prs, s, deck.contact_email)
            continue
        renderer = LAYOUT_RENDERERS.get(s.layout)
        if renderer is None:
            raise ValueError(f"Unknown layout: {s.layout!r}")
        renderer(prs, s)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _convert_pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> Path:
    """Convert a .pptx to .pdf using LibreOffice in headless mode.

    LibreOffice is the only converter that preserves embedded fonts, fills,
    and images with full fidelity to the brand template. If `soffice` is not
    on PATH, raise so the caller can either install it or fall back to PPTX.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        raise RuntimeError(
            "PDF output requires LibreOffice (soffice/libreoffice) on PATH. "
            "Install LibreOffice or rerun with --format pptx."
        )
    pdf_path = Path(pdf_path)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(pptx_path)],
            capture_output=True,
            text=True,
            check=False,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed (exit={result.returncode}): "
                f"{result.stderr.strip() or result.stdout.strip()}"
            )
        produced = Path(tmp) / (pptx_path.stem + ".pdf")
        if not produced.exists():
            raise RuntimeError(f"Conversion ran but no PDF was produced at {produced}")
        shutil.move(str(produced), str(pdf_path))
    return pdf_path


def build_deck(deck: Deck, output_path: str | Path, fmt: str = "pptx") -> Path:
    """Render a Deck to disk in the requested format and return the output path.

    Parameters
    ----------
    deck : Deck
        The structured deck definition.
    output_path : str | Path
        Destination file. The extension is normalized to match `fmt`.
    fmt : {"pptx", "pdf"}
        Output format. ``"pptx"`` writes an editable PowerPoint file.
        ``"pdf"`` first writes a temporary PPTX, then runs LibreOffice to
        produce a PDF. Defaults to ``"pptx"``.
    """
    fmt = fmt.lower()
    if fmt not in ("pptx", "pdf"):
        raise ValueError(f"format must be 'pptx' or 'pdf', got {fmt!r}")
    output_path = Path(output_path).with_suffix("." + fmt)
    if fmt == "pptx":
        return _render_pptx(deck, output_path)
    # PDF: write to a temp pptx first, then convert.
    with tempfile.TemporaryDirectory() as tmp:
        tmp_pptx = Path(tmp) / (output_path.stem + ".pptx")
        _render_pptx(deck, tmp_pptx)
        return _convert_pptx_to_pdf(tmp_pptx, output_path)



def build_from_json(json_path: str | Path, output_path: str | Path, fmt: str = "pptx") -> Path:
    """Build a deck from a JSON file matching the Deck schema."""
    data = json.loads(Path(json_path).read_text(encoding="utf-8"))
    deck = Deck(
        strategy_name=data["strategy_name"],
        contact_email=data.get("contact_email", "research@kaxanuk.mx"),
        slides=[
            Slide(
                layout=s["layout"],
                title=s.get("title"),
                subtitle=s.get("subtitle"),
                body=s.get("body", []),
                image_path=s.get("image_path"),
                page_number=s.get("page_number"),
            )
            for s in data["slides"]
        ],
    )
    return build_deck(deck, output_path, fmt=fmt)


# ---------------------------------------------------------------------------
# Example / smoke test.
# ---------------------------------------------------------------------------

def _example_deck() -> Deck:
    name = "Liquidity-Weighted Trend Strategy"
    return Deck(
        strategy_name=name,
        slides=[
            Slide(layout="cover"),
            Slide(layout="title", title="Strategy Presentation", subtitle=name),
            Slide(
                layout="content",
                title="Content",
                page_number=3,
                body=[
                    {
                        "kind": "arrows",
                        "items": ["Strategy Definition", "Feature Engineering", "Strategy Design"],
                    },
                ],
            ),
            Slide(layout="divider", title="Strategy Definition", subtitle=name),
            Slide(
                layout="content",
                title="Idea Description",
                page_number=5,
                body=[
                    {"kind": "callout", "text": "This strategy explores trend persistence in U.S. equities using a simple but robust signal."},
                    {"kind": "para", "text": "The core idea is that stocks experiencing sustained upward price trends tend to continue outperforming due to:"},
                    {"kind": "bullets", "items": ["Gradual information diffusion", "Institutional trading flows", "Behavioral underreaction"]},
                    {"kind": "para", "text": "Trend persistence is identified through a moving average crossover signal."},
                    {"kind": "para", "text": "Stocks are considered bullish when:"},
                    {"kind": "bullets", "items": ["SMA(50) > SMA(200)"]},
                ],
            ),
            Slide(layout="contact", title="Contact"),
        ],
    )


def _main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--example", action="store_true", help="Render a small sample deck for QA.")
    parser.add_argument("--input", help="Path to a Deck JSON file.")
    parser.add_argument("--output", required=True, help="Destination file path. Extension is normalized to match --format.")
    parser.add_argument(
        "--format",
        choices=("pptx", "pdf"),
        default="pptx",
        help="Output format. 'pdf' requires LibreOffice on PATH.",
    )
    args = parser.parse_args()
    if args.example:
        path = build_deck(_example_deck(), args.output, fmt=args.format)
    elif args.input:
        path = build_from_json(args.input, args.output, fmt=args.format)
    else:
        parser.error("provide either --example or --input")
    print(f"Wrote {path}")


if __name__ == "__main__":
    _main()
